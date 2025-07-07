#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
多格式输出转换器
支持将文档转换为多种格式：PDF、Word、PowerPoint等
"""

import os
import json
import base64
from pathlib import Path
from typing import Dict, List, Any, Optional
import re
from datetime import datetime

try:
    import markdown
    from weasyprint import HTML, CSS
    WEASYPRINT_AVAILABLE = True
except ImportError:
    WEASYPRINT_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class FormatConverter:
    """多格式转换器"""
    
    def __init__(self):
        self.supported_formats = self._get_supported_formats()
        
    def _get_supported_formats(self) -> List[str]:
        """获取支持的输出格式"""
        formats = ["HTML", "MARKDOWN"]
        
        if WEASYPRINT_AVAILABLE:
            formats.append("PDF")
        if DOCX_AVAILABLE:
            formats.append("DOCX")
        if PPTX_AVAILABLE:
            formats.append("PPTX")
            
        return formats
    
    def convert_to_format(self, content: str, target_format: str, 
                         styles: Dict[str, Any] = None, 
                         metadata: Dict[str, str] = None) -> Dict[str, Any]:
        """转换内容到指定格式"""
        target_format = target_format.upper()
        
        if target_format not in self.supported_formats:
            return {
                "success": False,
                "error": f"不支持的格式: {target_format}",
                "supported_formats": self.supported_formats
            }
        
        try:
            if target_format == "HTML":
                return self._convert_to_html(content, styles, metadata)
            elif target_format == "PDF":
                return self._convert_to_pdf(content, styles, metadata)
            elif target_format == "DOCX":
                return self._convert_to_docx(content, styles, metadata)
            elif target_format == "PPTX":
                return self._convert_to_pptx(content, styles, metadata)
            elif target_format == "MARKDOWN":
                return self._convert_to_markdown(content, styles, metadata)
            else:
                return {
                    "success": False,
                    "error": f"格式 {target_format} 的转换器尚未实现"
                }
        except Exception as e:
            return {
                "success": False,
                "error": f"转换过程中发生错误: {str(e)}"
            }
    
    def _convert_to_html(self, content: str, styles: Dict[str, Any] = None, 
                        metadata: Dict[str, str] = None) -> Dict[str, Any]:
        """转换为HTML格式"""
        # 将Markdown转换为HTML
        html_content = markdown.markdown(content, extensions=['extra', 'codehilite', 'toc'])
        
        # 生成CSS样式
        css_styles = self._generate_css(styles or {})
        
        # 构建完整的HTML文档
        full_html = f"""
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{metadata.get('title', '文档') if metadata else '文档'}</title>
    <style>
        {css_styles}
    </style>
</head>
<body>
    <div class="document-container">
        {html_content}
    </div>
    
    <footer class="document-footer">
        <p>生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
        {f'<p>作者: {metadata.get("author", "")}</p>' if metadata and metadata.get("author") else ''}
    </footer>
</body>
</html>
"""
        
        # 保存文件
        output_path = self._save_file(full_html, "html")
        
        return {
            "success": True,
            "format": "HTML",
            "file_path": output_path,
            "content": full_html,
            "size": len(full_html)
        }
    
    def _convert_to_pdf(self, content: str, styles: Dict[str, Any] = None, 
                       metadata: Dict[str, str] = None) -> Dict[str, Any]:
        """转换为PDF格式"""
        if not WEASYPRINT_AVAILABLE:
            return {
                "success": False,
                "error": "PDF转换需要安装weasyprint库: pip install weasyprint"
            }
        
        # 先转换为HTML
        html_result = self._convert_to_html(content, styles, metadata)
        if not html_result["success"]:
            return html_result
        
        html_content = html_result["content"]
        
        # 创建PDF特定的CSS
        pdf_css = self._generate_pdf_css(styles or {})
        
        try:
            # 使用WeasyPrint转换为PDF
            html_doc = HTML(string=html_content)
            css_doc = CSS(string=pdf_css)
            
            output_path = self._get_output_path("pdf")
            html_doc.write_pdf(output_path, stylesheets=[css_doc])
            
            # 获取文件大小
            file_size = os.path.getsize(output_path)
            
            return {
                "success": True,
                "format": "PDF",
                "file_path": output_path,
                "size": file_size
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"PDF生成失败: {str(e)}"
            }
    
    def _convert_to_docx(self, content: str, styles: Dict[str, Any] = None, 
                        metadata: Dict[str, str] = None) -> Dict[str, Any]:
        """转换为Word文档格式"""
        if not DOCX_AVAILABLE:
            return {
                "success": False,
                "error": "Word转换需要安装python-docx库: pip install python-docx"
            }
        
        try:
            # 创建新的Word文档
            doc = Document()
            
            # 设置文档属性
            if metadata:
                core_props = doc.core_properties
                core_props.title = metadata.get('title', '')
                core_props.author = metadata.get('author', '')
                core_props.subject = metadata.get('subject', '')
                core_props.comments = metadata.get('description', '')
            
            # 解析Markdown内容
            lines = content.split('\n')
            
            for line in lines:
                line = line.strip()
                
                if not line:
                    # 空行
                    continue
                
                # 处理标题
                if line.startswith('#'):
                    level = min(len(line) - len(line.lstrip('#')), 6)
                    title_text = line.lstrip('#').strip()
                    
                    heading = doc.add_heading(title_text, level=level)
                    self._apply_heading_style(heading, level, styles)
                
                # 处理图片
                elif line.startswith('!['):
                    # 简单的图片占位符
                    doc.add_paragraph(f"[图片: {line}]")
                
                # 处理列表
                elif line.startswith(('-', '*', '+')):
                    list_text = line[1:].strip()
                    p = doc.add_paragraph(list_text, style='List Bullet')
                
                elif re.match(r'^\d+\.', line):
                    list_text = re.sub(r'^\d+\.\s*', '', line)
                    p = doc.add_paragraph(list_text, style='List Number')
                
                # 处理普通段落
                else:
                    if line:
                        p = doc.add_paragraph(line)
                        self._apply_paragraph_style(p, styles)
            
            # 保存文档
            output_path = self._get_output_path("docx")
            doc.save(output_path)
            
            file_size = os.path.getsize(output_path)
            
            return {
                "success": True,
                "format": "DOCX",
                "file_path": output_path,
                "size": file_size
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Word文档生成失败: {str(e)}"
            }
    
    def _convert_to_pptx(self, content: str, styles: Dict[str, Any] = None, 
                        metadata: Dict[str, str] = None) -> Dict[str, Any]:
        """转换为PowerPoint演示文稿格式"""
        if not PPTX_AVAILABLE:
            return {
                "success": False,
                "error": "PowerPoint转换需要安装python-pptx库: pip install python-pptx"
            }
        
        try:
            # 创建新的PowerPoint演示文稿
            prs = Presentation()
            
            # 解析内容为幻灯片
            slides_content = self._parse_content_for_slides(content)
            
            for slide_data in slides_content:
                # 添加幻灯片
                slide_layout = prs.slide_layouts[1]  # 标题和内容布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 设置标题
                title = slide.shapes.title
                title.text = slide_data['title']
                
                # 设置内容
                content_placeholder = slide.placeholders[1]
                tf = content_placeholder.text_frame
                
                for item in slide_data['content']:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
            
            # 如果没有内容，添加标题页
            if not slides_content:
                slide_layout = prs.slide_layouts[0]  # 标题布局
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = metadata.get('title', '文档标题') if metadata else '文档标题'
                subtitle.text = metadata.get('author', '') if metadata else ''
            
            # 保存演示文稿
            output_path = self._get_output_path("pptx")
            prs.save(output_path)
            
            file_size = os.path.getsize(output_path)
            
            return {
                "success": True,
                "format": "PPTX",
                "file_path": output_path,
                "size": file_size
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"PowerPoint生成失败: {str(e)}"
            }
    
    def _convert_to_markdown(self, content: str, styles: Dict[str, Any] = None, 
                           metadata: Dict[str, str] = None) -> Dict[str, Any]:
        """转换为优化的Markdown格式"""
        # 添加元数据
        enhanced_content = ""
        
        if metadata:
            enhanced_content += "---\n"
            for key, value in metadata.items():
                enhanced_content += f"{key}: {value}\n"
            enhanced_content += "---\n\n"
        
        enhanced_content += content
        
        # 添加生成信息
        enhanced_content += f"\n\n---\n*文档生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n"
        
        # 保存文件
        output_path = self._save_file(enhanced_content, "md")
        
        return {
            "success": True,
            "format": "MARKDOWN",
            "file_path": output_path,
            "content": enhanced_content,
            "size": len(enhanced_content)
        }
    
    def _parse_content_for_slides(self, content: str) -> List[Dict[str, Any]]:
        """解析内容为幻灯片结构"""
        lines = content.split('\n')
        slides = []
        current_slide = None
        
        for line in lines:
            line = line.strip()
            
            # 一级标题作为新幻灯片
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                
                current_slide = {
                    'title': line[2:].strip(),
                    'content': []
                }
            
            # 二级标题作为新幻灯片
            elif line.startswith('## '):
                if current_slide:
                    slides.append(current_slide)
                
                current_slide = {
                    'title': line[3:].strip(),
                    'content': []
                }
            
            # 其他内容添加到当前幻灯片
            elif line and current_slide:
                # 跳过图片引用
                if not line.startswith('!['):
                    # 清理Markdown格式
                    clean_line = re.sub(r'[*_`]', '', line)
                    if clean_line.strip():
                        current_slide['content'].append(clean_line)
        
        # 添加最后一个幻灯片
        if current_slide:
            slides.append(current_slide)
        
        return slides
    
    def _generate_css(self, styles: Dict[str, Any]) -> str:
        """生成CSS样式"""
        typography = styles.get('typography', {})
        colors = styles.get('colors', {})
        spacing = styles.get('spacing', {})
        image_design = styles.get('image_design', {})
        
        css = f"""
        .document-container {{
            max-width: 900px;
            margin: 0 auto;
            padding: 40px 20px;
            font-family: {typography.get('primary_font', "'Segoe UI', 'Microsoft YaHei', sans-serif")};
            line-height: {typography.get('line_height', '1.6')};
            color: {colors.get('text', '#333333')};
            background-color: {colors.get('background', '#ffffff')};
        }}
        
        h1, h2, h3, h4, h5, h6 {{
            color: {colors.get('primary', '#2c3e50')};
            margin-top: {spacing.get('title_margin', '2em')};
            margin-bottom: 1em;
            font-weight: 600;
        }}
        
        h1 {{ font-size: {typography.get('heading_sizes', {}).get('h1', '2.5em')}; }}
        h2 {{ font-size: {typography.get('heading_sizes', {}).get('h2', '2em')}; }}
        h3 {{ font-size: {typography.get('heading_sizes', {}).get('h3', '1.5em')}; }}
        
        p {{
            margin-bottom: {spacing.get('paragraph_margin', '1em')};
            text-align: justify;
        }}
        
        img {{
            max-width: {image_design.get('max_width', '100%')};
            height: auto;
            border-radius: {image_design.get('border_radius', '8px')};
            box-shadow: {image_design.get('shadow', '0 4px 8px rgba(0,0,0,0.1)')};
            margin: 1em 0;
            display: block;
            margin-left: auto;
            margin-right: auto;
        }}
        
        code {{
            background-color: #f4f4f4;
            padding: 2px 4px;
            border-radius: 3px;
            font-family: 'Consolas', 'Monaco', monospace;
        }}
        
        pre {{
            background-color: #f8f8f8;
            padding: 15px;
            border-radius: 5px;
            overflow-x: auto;
            border-left: 4px solid {colors.get('primary', '#2c3e50')};
        }}
        
        ul, ol {{
            margin-bottom: 1em;
            padding-left: 2em;
        }}
        
        li {{
            margin-bottom: 0.5em;
        }}
        
        blockquote {{
            border-left: 4px solid {colors.get('secondary', '#95a5a6')};
            margin: 1em 0;
            padding-left: 1em;
            color: #666;
            font-style: italic;
        }}
        
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 1em 0;
        }}
        
        th, td {{
            border: 1px solid #ddd;
            padding: 8px 12px;
            text-align: left;
        }}
        
        th {{
            background-color: {colors.get('primary', '#2c3e50')};
            color: white;
        }}
        
        .document-footer {{
            margin-top: 3em;
            padding-top: 2em;
            border-top: 1px solid #eee;
            text-align: center;
            color: #888;
            font-size: 0.9em;
        }}
        """
        
        return css
    
    def _generate_pdf_css(self, styles: Dict[str, Any]) -> str:
        """生成PDF专用CSS"""
        base_css = self._generate_css(styles)
        
        pdf_specific = """
        @page {
            margin: 2cm;
            size: A4;
        }
        
        .document-container {
            max-width: none;
            margin: 0;
            padding: 0;
        }
        
        h1, h2, h3 {
            page-break-after: avoid;
        }
        
        img {
            page-break-inside: avoid;
            max-height: 20cm;
        }
        
        table {
            page-break-inside: avoid;
        }
        """
        
        return base_css + pdf_specific
    
    def _apply_heading_style(self, heading, level: int, styles: Dict[str, Any]):
        """应用标题样式到Word文档"""
        if not styles:
            return
        
        colors = styles.get('colors', {})
        typography = styles.get('typography', {})
        
        # 设置颜色
        primary_color = colors.get('primary', '#2c3e50')
        if primary_color.startswith('#'):
            # 简单的颜色转换（实际应用中可能需要更复杂的转换）
            pass
    
    def _apply_paragraph_style(self, paragraph, styles: Dict[str, Any]):
        """应用段落样式到Word文档"""
        if not styles:
            return
        
        spacing = styles.get('spacing', {})
        
        # 设置段落间距
        if 'paragraph_margin' in spacing:
            # 这里可以设置段落格式
            pass
    
    def _get_output_path(self, extension: str) -> str:
        """获取输出文件路径"""
        output_dir = Path("output")
        output_dir.mkdir(exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"formatted_document_{timestamp}.{extension}"
        
        return str(output_dir / filename)
    
    def _save_file(self, content: str, extension: str) -> str:
        """保存文件"""
        output_path = self._get_output_path(extension)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return output_path
    
    def get_available_formats(self) -> Dict[str, Dict[str, Any]]:
        """获取可用格式及其描述"""
        formats_info = {
            "HTML": {
                "name": "HTML网页",
                "description": "适合在线查看和分享",
                "available": True,
                "extensions": [".html", ".htm"]
            },
            "MARKDOWN": {
                "name": "Markdown文档",
                "description": "轻量级标记语言，便于编辑",
                "available": True,
                "extensions": [".md", ".markdown"]
            },
            "PDF": {
                "name": "PDF文档",
                "description": "适合打印和正式分发",
                "available": WEASYPRINT_AVAILABLE,
                "extensions": [".pdf"],
                "requirements": "pip install weasyprint" if not WEASYPRINT_AVAILABLE else None
            },
            "DOCX": {
                "name": "Word文档",
                "description": "Microsoft Word格式，便于编辑",
                "available": DOCX_AVAILABLE,
                "extensions": [".docx"],
                "requirements": "pip install python-docx" if not DOCX_AVAILABLE else None
            },
            "PPTX": {
                "name": "PowerPoint演示文稿",
                "description": "演示文稿格式，适合展示",
                "available": PPTX_AVAILABLE,
                "extensions": [".pptx"],
                "requirements": "pip install python-pptx" if not PPTX_AVAILABLE else None
            }
        }
        
        return formats_info

# 便捷函数
def convert_document(content: str, target_format: str, styles: Dict[str, Any] = None, 
                    metadata: Dict[str, str] = None) -> Dict[str, Any]:
    """转换文档的便捷函数"""
    converter = FormatConverter()
    return converter.convert_to_format(content, target_format, styles, metadata)

def get_supported_formats() -> List[str]:
    """获取支持的格式列表"""
    converter = FormatConverter()
    return converter.supported_formats

if __name__ == "__main__":
    # 测试格式转换功能
    test_content = """
# 测试文档

## 介绍

这是一个测试文档，用于验证多格式转换功能。

### 特性

- 支持多种输出格式
- 自动样式应用
- 高质量转换

![测试图片](test.jpg)

## 结论

格式转换功能正常工作。
"""
    
    converter = FormatConverter()
    
    print("支持的格式:", converter.supported_formats)
    
    # 测试HTML转换
    result = converter.convert_to_format(test_content, "HTML")
    print("HTML转换结果:", result["success"])
    
    # 获取格式信息
    formats_info = converter.get_available_formats()
    for fmt, info in formats_info.items():
        status = "✅" if info["available"] else "❌"
        print(f"{status} {fmt}: {info['description']}")
        if info.get("requirements"):
            print(f"   需要: {info['requirements']}")